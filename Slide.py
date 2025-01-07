from pptx import Presentation
from pptx.util import Inches, Pt

# Créer une présentation
presentation = Presentation()

# Slide 1 - Page de titre
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "British Airways Customer Review Analysis"
subtitle.text = "Data Insights and Recommendations"

# Slide 2 - Graphique de distribution des sentiments
slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide2.shapes.title
title.text = "Sentiment Distribution"
left = Inches(1)
top = Inches(1.5)
slide2.shapes.add_picture('image 3.png', left, top, width=Inches(7))  # Remplacez par le chemin de votre image

# Slide 3 - Répartition des sentiments en chiffres
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide3.shapes.title
title.text = "Sentiment Count Overview"
content = slide3.placeholders[1]
content.text = (
    "Répartition des sentiments dans les avis :\n\n"
    "- **Positive** : 2,701 avis\n"
    "- **Negative** : 1,172 avis\n"
    "- **Neutral** : 21 avis\n"
)

# Slide 4 - Topics principaux (LDA)
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide4.shapes.title
title.text = "Top Topics in Reviews (LDA)"
content = slide4.placeholders[1]
content.text = (
    "Topics principaux identifiés dans les avis :\n\n"
    "- **Topic 1** : flight, meal, food, chicken, british\n"
    "- **Topic 2** : flight, crew, good, ba, time\n"
    "- **Topic 3** : seat, seats, economy, flight, food\n"
    "- **Topic 4** : flight, ba, british, airways, london\n"
    "- **Topic 5** : ba, flight, service, class, british\n"
)

# Slide 5 - Insights clés
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide5.shapes.title
title.text = "Key Insights"
content = slide5.placeholders[1]
content.text = (
    "- **Low Ratings** : 70% des avis ont une note de 1 ou 2.\n"
    "- **Negative Sentiment** : 65% des avis sont négatifs.\n"
    "- **Complaints** : Retards, mauvais service client, problèmes de bagages.\n"
    "- **Positive Feedback** : Confort des sièges et personnel aimable.\n"
)

# Slide 6 - Recommandations
slide6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide6.shapes.title
title.text = "Recommendations"
content = slide6.placeholders[1]
content.text = (
    "1. Renforcer la formation du personnel pour un meilleur service client.\n"
    "2. Résoudre les problèmes liés aux retards et à la gestion des bagages.\n"
    "3. Améliorer la communication avec les passagers lors des perturbations.\n"
)

# Sauvegarder la présentation
presentation.save('British_Airways_Analysis_Updated.pptx')
print("La présentation a été créée : 'British_Airways_Analysis_Updated1.pptx'.")
